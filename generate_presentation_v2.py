"""Equb Presentation v2 — Visual redesign with real app screenshots."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
import os, sys

BASE = r"c:\Users\tnega\OneDrive\Desktop\እቁብ (Equb)"

# ── Colors ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x1D, 0x35)
NAVY2  = RGBColor(0x16, 0x2E, 0x4E)
GOLD   = RGBColor(0xF5, 0xA6, 0x23)
GREEN  = RGBColor(0x07, 0x8A, 0x3C)
RED    = RGBColor(0xDC, 0x14, 0x3C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF4, 0xF6, 0xFA)
MGRAY  = RGBColor(0xD8, 0xE2, 0xEF)
SLATE  = RGBColor(0x4A, 0x5E, 0x78)
DSLATE = RGBColor(0x2D, 0x3A, 0x4E)
BLUE   = RGBColor(0x1A, 0x73, 0xE8)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)
CYAN   = RGBColor(0x00, 0xB8, 0xD4)
LBLUE  = RGBColor(0xE8, 0xF0, 0xFE)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # fully blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill: RGBColor):
    shp = slide.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp

def oval(slide, x, y, w, h, fill: RGBColor):
    shp = slide.shapes.add_shape(9, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp

def txt(slide, text, x, y, w, h, size, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.size   = Pt(size)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb

def grad(shape, hex1, hex2, angle_deg=270):
    """Apply a linear gradient fill to a shape (rectangles / ovals only)."""
    ang = angle_deg * 60000
    sp  = shape._element
    spPr = sp.spPr
    for child in list(spPr):
        tag = child.tag
        if tag in (qn('a:solidFill'), qn('a:gradFill'),
                   qn('a:noFill'), qn('a:pattFill'), qn('a:blipFill')):
            spPr.remove(child)
    gfill = parse_xml(
        f'<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{hex1}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{hex2}"/></a:gs>'
        f'</a:gsLst>'
        f'<a:lin ang="{ang}" scaled="0"/>'
        f'</a:gradFill>'
    )
    spPr.insert(0, gfill)

def photo(slide, filename, x, y, w, h=None):
    path = os.path.join(BASE, filename)
    if h:
        return slide.shapes.add_picture(path, x, y, w, h)
    return slide.shapes.add_picture(path, x, y, w)

def header_band(slide, title, subtitle=None):
    bg = rect(slide, 0, 0, W, Inches(1.15), NAVY)
    grad(bg, '060F1E', '162E4E', 180)
    rect(slide, 0, Inches(1.15), W, Inches(0.055), GOLD)
    txt(slide, title, Inches(0.55), Inches(0.17), Inches(11), Inches(0.72),
        30, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.55), Inches(0.83), Inches(11), Inches(0.34),
            12, color=RGBColor(0x88, 0xA4, 0xCC))

def flag_stripe_left(slide):
    rect(slide, 0, 0,        Inches(0.16), H * 0.33, GREEN)
    rect(slide, 0, H * 0.33, Inches(0.16), H * 0.34, GOLD)
    rect(slide, 0, H * 0.67, Inches(0.16), H * 0.33, RED)

def flag_stripe_right(slide):
    x = W - Inches(0.16)
    rect(slide, x, 0,        Inches(0.16), H * 0.33, GREEN)
    rect(slide, x, H * 0.33, Inches(0.16), H * 0.34, GOLD)
    rect(slide, x, H * 0.67, Inches(0.16), H * 0.33, RED)


# ── SLIDE 1 — TITLE ──────────────────────────────────────────────────────────

def slide_title():
    slide = prs.slides.add_slide(blank)

    bg = rect(slide, 0, 0, W, H, NAVY)
    grad(bg, '050D1A', '102340', 145)

    flag_stripe_left(slide)

    # decorative rings (right side background)
    for sz, cx, cy in [(Inches(5.8), Inches(9.8), Inches(-1.5)),
                        (Inches(3.5), Inches(11.5), Inches(4.8))]:
        o = slide.shapes.add_shape(9, cx, cy, sz, sz)
        o.fill.background()
        o.line.color.rgb = RGBColor(0x18, 0x34, 0x58)
        o.line.width = Pt(1.8)

    # right panel: screenshot
    right_bg = rect(slide, Inches(7.9), 0, Inches(5.43), H, RGBColor(0x0A, 0x1C, 0x32))
    grad(right_bg, '08162A', '0E2540', 180)

    photo(slide, "check_superadmin.png",
          Inches(8.05), Inches(0.5), Inches(5.1), Inches(6.1))

    # bottom fade on screenshot
    fade = rect(slide, Inches(7.9), Inches(5.5), Inches(5.43), Inches(2.0),
                RGBColor(0x08, 0x16, 0x2A))
    grad(fade, '08162A00', '08162A', 270)

    # left content
    # Gold circle icon
    icon = oval(slide, Inches(0.55), Inches(0.75), Inches(1.35), Inches(1.35), GOLD)
    txt(slide, "እ", Inches(0.55), Inches(0.82), Inches(1.35), Inches(1.2),
        52, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    txt(slide, "እቁብ",
        Inches(0.55), Inches(2.1), Inches(7.0), Inches(1.7),
        90, bold=True, color=WHITE)

    txt(slide, "Equb Management Platform",
        Inches(0.58), Inches(3.7), Inches(7.0), Inches(0.65),
        26, bold=False, color=GOLD)

    rect(slide, Inches(0.58), Inches(4.45), Inches(4.8), Inches(0.05), GOLD)

    txt(slide, "Digital infrastructure for Ethiopia's\nrotating savings & credit groups",
        Inches(0.58), Inches(4.6), Inches(7.0), Inches(0.9),
        14, color=RGBColor(0x7A, 0x9E, 0xC0))

    # stat counter tiles
    stats = [
        ("113+", "Member Spots",    GREEN),
        ("6",    "SMS Templates",   GOLD),
        ("3",    "Access Roles",    BLUE),
        ("v3.1", "Live on Railway", RED),
    ]
    sx = Inches(0.55)
    for i, (val, label, color) in enumerate(stats):
        bx = sx + i * Inches(1.82)
        tile = rect(slide, bx, Inches(5.65), Inches(1.68), Inches(1.45),
                    RGBColor(0x0C, 0x20, 0x38))
        rect(slide, bx, Inches(5.65), Inches(1.68), Inches(0.065), color)
        txt(slide, val,   bx, Inches(5.73), Inches(1.68), Inches(0.72),
            34, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(slide, label, bx, Inches(6.42), Inches(1.68), Inches(0.5),
            10, color=RGBColor(0x7A, 0x9E, 0xC0), align=PP_ALIGN.CENTER)


# ── SLIDE 2 — HOW EQUB WORKS ─────────────────────────────────────────────────

def slide_concept():
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, W, H, LGRAY)
    header_band(slide, "How Equb Works",
                "A centuries-old Ethiopian savings tradition — now digital and transparent")

    steps = [
        ("1", "Members\nJoin",         "Each member buys a\nspot in the cycle",    GREEN),
        ("2", "Weekly\nContributions",  "Every member pays\ntheir weekly share",    GOLD),
        ("3", "Draw\nHeld",            "A winner is drawn\nfrom open spots",        BLUE),
        ("4", "Winner\nPaid",          "Full pot disbursed\nby cheque",             RED),
    ]

    circ_r  = Inches(1.45)
    step_w  = Inches(2.8)
    gap     = Inches(0.55)
    total_w = 4 * step_w + 3 * gap
    sx      = (W - total_w) / 2
    cy      = Inches(3.0)

    for i, (num, title, desc, color) in enumerate(steps):
        x = sx + i * (step_w + gap)
        cx = x + (step_w - circ_r) / 2

        # Shadow circle (offset slightly)
        shadow = oval(slide, cx + Inches(0.07), cy - circ_r / 2 + Inches(0.07),
                      circ_r, circ_r, MGRAY)

        # Main circle
        c = oval(slide, cx, cy - circ_r / 2, circ_r, circ_r, color)
        grad(c,
             {GREEN: '078A3C', GOLD: 'F5A623', BLUE: '1A73E8', RED: 'DC143C'}[color],
             {GREEN: '056A2E', GOLD: 'D4880E', BLUE: '0F5CC0', RED: 'A80E2A'}[color],
             135)

        # Number
        txt(slide, num, cx, cy - circ_r / 2 + Inches(0.18),
            circ_r, circ_r - Inches(0.3), 60, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

        # Title
        txt(slide, title, x, cy + circ_r / 2 + Inches(0.18),
            step_w, Inches(0.7), 14, bold=True,
            color=DSLATE, align=PP_ALIGN.CENTER)

        # Description
        txt(slide, desc, x, cy + circ_r / 2 + Inches(0.92),
            step_w, Inches(0.75), 11, color=SLATE, align=PP_ALIGN.CENTER)

        # Arrow connector
        if i < 3:
            ax = cx + circ_r + Inches(0.06)
            aw = gap - Inches(0.12)
            arrow = rect(slide, ax, cy - Inches(0.025), aw - Inches(0.18), Inches(0.05), MGRAY)
            # Arrowhead: right-pointing arrow using a rotated shape
            tip = slide.shapes.add_shape(13, ax + aw - Inches(0.22),
                                          cy - Inches(0.12), Inches(0.22), Inches(0.24))
            tip.fill.solid(); tip.fill.fore_color.rgb = MGRAY
            tip.line.fill.background()

    # Bottom banner
    banner = rect(slide, 0, Inches(6.35), W, Inches(1.15), NAVY)
    grad(banner, '0B1D35', '162E4E', 180)
    txt(slide,
        "Every member wins exactly once per cycle — guaranteed, transparent, and recorded",
        Inches(0.8), Inches(6.57), Inches(11.73), Inches(0.65),
        18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ── SLIDE 3 — BY THE NUMBERS ─────────────────────────────────────────────────

def slide_numbers():
    slide = prs.slides.add_slide(blank)
    bg = rect(slide, 0, 0, W, H, NAVY)
    grad(bg, '050D1A', '0F2845', 135)

    rect(slide, 0, Inches(1.15), W, Inches(0.055), GOLD)
    txt(slide, "By the Numbers", Inches(0.55), Inches(0.2),
        Inches(10), Inches(0.78), 32, bold=True, color=WHITE)
    txt(slide, "The scale of the Equb platform at a glance",
        Inches(0.55), Inches(0.84), Inches(10), Inches(0.36),
        13, color=RGBColor(0x88, 0xA4, 0xCC))

    nums = [
        ("113",  "Member Spots\nper Cycle",                GREEN,  "+5 Association spots"),
        ("118",  "Total Spots\n(incl. Association)",       GOLD,   "Full + half-share splits"),
        ("6",    "SMS Templates\nfor Automation",          BLUE,   "Powered by Africa's Talking"),
        ("3",    "User Roles\n(Superadmin/Admin/Cashier)", RED,    "Granular permission matrix"),
        ("48h",  "Before Draw\nReminder Sent",             PURPLE, "Auto-scheduled by APScheduler"),
        ("5",    "Failed Logins\nbefore Lockout",          CYAN,   "5-minute rate-limit ban"),
    ]

    cw = Inches(4.06)
    ch = Inches(2.38)
    gx = Inches(0.16)
    gy = Inches(0.18)
    sx = Inches(0.38)
    sy = Inches(1.35)

    for i, (val, label, color, sub) in enumerate(nums):
        row, col = divmod(i, 3)
        x = sx + col * (cw + gx)
        y = sy + row * (ch + gy)

        card = rect(slide, x, y, cw, ch, RGBColor(0x0D, 0x21, 0x3C))
        grad(card, '0D213C', '172E4E', 145)
        rect(slide, x, y, cw, Inches(0.065), color)

        txt(slide, val, x, y + Inches(0.1), cw, Inches(1.1),
            68, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(slide, label, x, y + Inches(1.12), cw, Inches(0.8),
            12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, sub, x, y + Inches(1.9), cw, Inches(0.38),
            10, color=RGBColor(0x6A, 0x8A, 0xB0), align=PP_ALIGN.CENTER, italic=True)


# ── SLIDE 4 — DASHBOARD SCREENSHOT ───────────────────────────────────────────

def slide_dashboard():
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, W, H, LGRAY)
    header_band(slide, "The Platform — Dashboard",
                "Role-based interface for Superadmin, Admin, and Cashier")

    # Big screenshot right
    photo(slide, "check_superadmin.png",
          Inches(5.3), Inches(1.28), Inches(7.75), Inches(5.95))

    # callout dots + lines
    callouts = [
        (GREEN,  Inches(1.55), "Sidebar Navigation",
         "Dashboard · Members · Draws\nPayments · Reports · Notifications · Settings"),
        (GOLD,   Inches(2.75), "Ethiopian Flag Branding",
         "Green / Gold / Red stripe always visible\nin the app header for every user"),
        (BLUE,   Inches(3.95), "Active Cycle Selector",
         "Switch cycles from any page\nwithout losing context"),
        (RED,    Inches(5.15), "User Role Badge",
         "Current role shown in sidebar footer\nAdmin / Cashier / Superadmin"),
    ]
    for color, cy, title, desc in callouts:
        dot = oval(slide, Inches(0.3), cy - Inches(0.16), Inches(0.32), Inches(0.32), color)
        # horizontal line to text
        rect(slide, Inches(0.62), cy, Inches(0.35), Inches(0.03), MGRAY)
        txt(slide, title, Inches(1.05), cy - Inches(0.22),
            Inches(4.0), Inches(0.4), 12, bold=True, color=DSLATE)
        txt(slide, desc,  Inches(1.05), cy + Inches(0.2),
            Inches(4.0), Inches(0.65), 10, color=SLATE)
        rect(slide, Inches(0.3), cy + Inches(0.9), Inches(4.75), Inches(0.02), MGRAY)


# ── SLIDE 5 — CORE FEATURES ──────────────────────────────────────────────────

def slide_features():
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, W, H, LGRAY)
    header_band(slide, "Core Features",
                "Everything a growing Equb association needs")

    features = [
        ("M", "Members",        GREEN,  "Track 113+ spots with full/half-share splitting, status flags, and member notes"),
        ("$", "Payments",       GOLD,   "Weekly payment entry with automatic late/missed transitions after 3 days"),
        ("D", "Draws",          BLUE,   "Scheduled weekly draws — hold logic, winner selection, spot-sale override"),
        ("C", "Disbursements",  RED,    "Cheque payouts with 3 required guarantors, voucher deductions, net calc"),
        ("R", "Reports",        PURPLE, "Financial summary, balance sheet, vouchers, general ledger, print-to-PDF"),
        ("S", "SMS Alerts",     CYAN,   "6 templates via Africa's Talking: reminders, winner, missed, pot sold…"),
    ]

    cw = Inches(4.12)
    ch = Inches(2.42)
    gx = Inches(0.16)
    gy = Inches(0.16)
    sx = Inches(0.37)
    sy = Inches(1.32)

    for i, (icon, name, color, desc) in enumerate(features):
        row, col = divmod(i, 3)
        x = sx + col * (cw + gx)
        y = sy + row * (ch + gy)

        rect(slide, x, y, cw, ch, WHITE)
        rect(slide, x, y, cw, Inches(0.065), color)

        # icon circle
        ir = Inches(0.72)
        ix = x + Inches(0.28)
        iy = y + Inches(0.22)
        ic = oval(slide, ix, iy, ir, ir, color)
        g_map = {GREEN:'078A3C', GOLD:'F5A623', BLUE:'1A73E8',
                 RED:'DC143C', PURPLE:'9B59B6', CYAN:'00B8D4'}
        g2    = {GREEN:'056A2E', GOLD:'D4880E', BLUE:'0F5CC0',
                 RED:'A80E2A', PURPLE:'7D3F9A', CYAN:'008FA3'}
        grad(ic, g_map[color], g2[color], 135)
        txt(slide, icon, ix, iy + Inches(0.1), ir, ir - Inches(0.15),
            24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # feature name
        txt(slide, name, x + Inches(1.12), iy + Inches(0.1),
            cw - Inches(1.25), Inches(0.42), 14, bold=True, color=DSLATE)

        rect(slide, x + Inches(0.28), y + Inches(1.12),
             cw - Inches(0.56), Inches(0.02), MGRAY)

        txt(slide, desc, x + Inches(0.28), y + Inches(1.2),
            cw - Inches(0.56), Inches(1.12), 10, color=SLATE)


# ── SLIDE 6 — REPORTS SCREENSHOT ─────────────────────────────────────────────

def slide_reports():
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, W, H, LGRAY)
    header_band(slide, "Reports & Analytics",
                "Real-time financial visibility with print-to-PDF export")

    # Left: screenshot
    photo(slide, "pg_reports.png",
          Inches(0.3), Inches(1.28), Inches(7.25), Inches(5.95))

    # Right: feature callouts
    items = [
        (GOLD,  "Financial Summary",
         "Total collected · Association fund\nService fee · Collection rate · Voucher pool"),
        (GREEN, "Balance Sheet",
         "Full accounting view of all ETB\nmovements across the cycle"),
        (BLUE,  "General Ledger",
         "Itemised week-by-week transaction\nbreakdown for audits"),
        (RED,   "Disbursement Table",
         "Winner, date, gross, deductions,\nnet cheque, guarantor names"),
    ]
    y = Inches(1.55)
    for color, title, desc in items:
        dot = oval(slide, Inches(7.85), y, Inches(0.33), Inches(0.33), color)
        txt(slide, title, Inches(8.32), y - Inches(0.06),
            Inches(4.75), Inches(0.4), 12, bold=True, color=DSLATE)
        txt(slide, desc, Inches(8.32), y + Inches(0.35),
            Inches(4.75), Inches(0.6), 10, color=SLATE)
        rect(slide, Inches(7.85), y + Inches(1.05), Inches(5.15), Inches(0.02), MGRAY)
        y += Inches(1.28)

    # Print note chip
    note = rect(slide, Inches(7.85), Inches(6.45), Inches(5.15), Inches(0.68), LBLUE)
    grad(note, 'E8F0FE', 'D2E3FC', 180)
    txt(slide, "Print Report → exports formatted financial PDF directly from browser",
        Inches(8.05), Inches(6.52), Inches(4.9), Inches(0.55), 10, color=BLUE)


# ── SLIDE 7 — SMS NOTIFICATIONS ──────────────────────────────────────────────

def slide_notifications():
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, W, H, LGRAY)
    header_band(slide, "SMS Notifications",
                "Automated member communication via Africa's Talking API")

    # Right: screenshot
    photo(slide, "pg_notifications.png",
          Inches(6.55), Inches(1.28), Inches(6.5), Inches(5.95))

    # Left: 6 template cards
    templates = [
        (GREEN,  "Payment Reminder",   "Auto-sent 48 hours before every weekly draw"),
        (RED,    "Missed Payment",      "All members with missed payments notified"),
        (GOLD,   "Draw Winner",         "Winner receives personal confirmation SMS"),
        (BLUE,   "Payment Confirmed",   "Receipt sent after cashier records payment"),
        (PURPLE, "Pot On Hold",         "Draw postponed — all members informed"),
        (CYAN,   "Pot Sold",            "Spot sale completed — buyer/seller notified"),
    ]
    y = Inches(1.35)
    for color, name, desc in templates:
        card = rect(slide, Inches(0.3), y, Inches(6.0), Inches(0.75), WHITE)
        rect(slide, Inches(0.3), y, Inches(0.07), Inches(0.75), color)
        txt(slide, name, Inches(0.52), y + Inches(0.08),
            Inches(5.4), Inches(0.35), 12, bold=True, color=DSLATE)
        txt(slide, desc, Inches(0.52), y + Inches(0.43),
            Inches(5.4), Inches(0.28), 10, color=SLATE)
        y += Inches(0.83)


# ── SLIDE 8 — TECH STACK ─────────────────────────────────────────────────────

def slide_tech():
    slide = prs.slides.add_slide(blank)
    bg = rect(slide, 0, 0, W, H, NAVY)
    grad(bg, '050D1A', '0F2845', 145)
    rect(slide, 0, Inches(1.15), W, Inches(0.055), GOLD)
    txt(slide, "Technology Stack", Inches(0.55), Inches(0.2),
        Inches(10), Inches(0.78), 30, bold=True, color=WHITE)
    txt(slide, "Production-grade Python stack, deployed on Railway.app",
        Inches(0.55), Inches(0.84), Inches(10), Inches(0.36),
        13, color=RGBColor(0x88, 0xA4, 0xCC))

    layers = [
        ("FRONTEND",  BLUE,   ["Jinja2 HTML templates", "Tailwind CSS", "Chart.js charts", "Vanilla JavaScript"]),
        ("BACKEND",   GREEN,  ["FastAPI  ·  Uvicorn", "APScheduler (cron jobs)", "SQLAlchemy ORM", "Africa's Talking SMS"]),
        ("DATABASE",  GOLD,   ["PostgreSQL (production)", "SQLite (local dev)", "Railway managed DB", "Alembic migrations"]),
        ("SECURITY",  RED,    ["PBKDF2 password hashing", "CSRF protection", "CSP + HTTPS headers", "Rate-limiting & sessions"]),
    ]

    cw = Inches(3.02)
    gx = Inches(0.17)
    total = 4 * cw + 3 * gx
    sx = (W - total) / 2
    sy = Inches(1.35)
    ch = Inches(5.7)

    for i, (label, color, items) in enumerate(layers):
        x = sx + i * (cw + gx)

        head = rect(slide, x, sy, cw, Inches(0.68), color)
        g_map = {BLUE:'1A73E8', GREEN:'078A3C', GOLD:'F5A623', RED:'DC143C'}
        g2    = {BLUE:'0F5CC0', GREEN:'056A2E', GOLD:'D4880E', RED:'A80E2A'}
        grad(head, g_map[color], g2[color], 135)
        txt(slide, label, x, sy + Inches(0.13), cw, Inches(0.44),
            13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        body = rect(slide, x, sy + Inches(0.68), cw, ch - Inches(0.68),
                    RGBColor(0x0D, 0x21, 0x3C))
        grad(body, '0D213C', '172E4E', 180)

        for j, item in enumerate(items):
            iy = sy + Inches(0.92) + j * Inches(1.17)
            chip = rect(slide, x + Inches(0.2), iy, cw - Inches(0.4), Inches(0.85),
                        RGBColor(0x14, 0x2A, 0x48))
            rect(slide, x + Inches(0.2), iy, Inches(0.06), Inches(0.85), color)
            txt(slide, item, x + Inches(0.38), iy + Inches(0.18),
                cw - Inches(0.62), Inches(0.5), 11, color=WHITE)

    # Railway badge bottom
    badge = rect(slide, 0, Inches(6.85), W, Inches(0.65), RGBColor(0x08, 0x16, 0x2A))
    txt(slide, "Deployed on Railway.app  |  Auto-deploy from Git  |  Managed PostgreSQL add-on",
        0, Inches(6.93), W, Inches(0.48), 12, color=GOLD, align=PP_ALIGN.CENTER)


# ── SLIDE 9 — CLOSING ─────────────────────────────────────────────────────────

def slide_closing():
    slide = prs.slides.add_slide(blank)
    bg = rect(slide, 0, 0, W, H, NAVY)
    grad(bg, '050D1A', '102340', 135)

    flag_stripe_left(slide)
    flag_stripe_right(slide)

    # decorative rings
    for sz, cx, cy in [(Inches(6), Inches(-1.5), Inches(-1.8)),
                        (Inches(3.8), Inches(11.2), Inches(4.8)),
                        (Inches(2), Inches(9.5), Inches(0.2))]:
        o = slide.shapes.add_shape(9, cx, cy, sz, sz)
        o.fill.background()
        o.line.color.rgb = RGBColor(0x16, 0x30, 0x52)
        o.line.width = Pt(2)

    txt(slide, "Built for",
        Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.85),
        40, color=RGBColor(0x7A, 0x9E, 0xC0), align=PP_ALIGN.CENTER)

    txt(slide, "Ethiopia's Financial Communities",
        Inches(0.5), Inches(1.85), Inches(12.33), Inches(1.2),
        56, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    rect(slide, Inches(3.5), Inches(3.15), Inches(6.33), Inches(0.06), GOLD)

    txt(slide, "Transparent  ·  Automated  ·  Scalable  ·  Secure",
        Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.6),
        18, color=RGBColor(0x88, 0xA4, 0xCC), align=PP_ALIGN.CENTER)

    # 4 pillars
    pillars = [("Transparent", GREEN), ("Automated", BLUE),
               ("Scalable",    GOLD),  ("Secure",    RED)]
    pw = Inches(2.8)
    pg = Inches(0.18)
    px = (W - (4 * pw + 3 * pg)) / 2
    py = Inches(4.1)
    g_map = {GREEN:'078A3C', BLUE:'1A73E8', GOLD:'F5A623', RED:'DC143C'}
    g2    = {GREEN:'056A2E', BLUE:'0F5CC0', GOLD:'D4880E', RED:'A80E2A'}
    for label, color in pillars:
        p = rect(slide, px, py, pw, Inches(0.75), color)
        grad(p, g_map[color], g2[color], 135)
        txt(slide, label, px, py + Inches(0.14), pw, Inches(0.48),
            17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        px += pw + pg

    # App info
    txt(slide, "Equb Management Platform  ·  v3.1  ·  FastAPI + PostgreSQL + Railway",
        Inches(0.5), Inches(6.85), Inches(12.33), Inches(0.45),
        11, color=RGBColor(0x44, 0x68, 0x90), align=PP_ALIGN.CENTER)


# ── Build ─────────────────────────────────────────────────────────────────────

slide_title()
slide_concept()
slide_numbers()
slide_dashboard()
slide_features()
slide_reports()
slide_notifications()
slide_tech()
slide_closing()

out = os.path.join(BASE, "Equb_Presentation_v2.pptx")
prs.save(out)
sys.stdout.buffer.write(b"Done: Equb_Presentation_v2.pptx\n")

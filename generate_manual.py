"""Equb Management Platform — Full User & Admin Manual.
Visual style matches the overview presentation (solid fills, white cards, left accent bars).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, sys

BASE = r"c:\Users\tnega\OneDrive\Desktop\እቁብ (Equb)"
OUT  = os.path.join(BASE, "Equb_Manual_v2.pptx")

# ── Brand colors (same as overview presentation) ────────────────────────────
NAVY   = RGBColor(0x0F, 0x24, 0x3E)
GOLD   = RGBColor(0xF5, 0xA6, 0x23)
GREEN  = RGBColor(0x07, 0x8A, 0x3C)
RED    = RGBColor(0xDC, 0x14, 0x3C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF4, 0xF6, 0xF9)
SLATE  = RGBColor(0x3A, 0x4A, 0x5C)
BLUE   = RGBColor(0x1A, 0x73, 0xE8)
PURPLE = RGBColor(0x7B, 0x2F, 0xBE)
CYAN   = RGBColor(0x00, 0x97, 0xA7)
DGRAY  = RGBColor(0x2D, 0x3A, 0x4E)
LGOLD  = RGBColor(0xFE, 0xF7, 0xE6)
LRED   = RGBColor(0xFD, 0xED, 0xED)
LGREEN = RGBColor(0xE6, 0xF4, 0xEA)
LBLUE  = RGBColor(0xE8, 0xF0, 0xFE)
MGRAY  = RGBColor(0xDC, 0xE3, 0xED)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ── Core helpers (v1 style) ─────────────────────────────────────────────────

def R(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def OV(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(9, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def T(slide, text, x, y, w, h, size, bold=False,
      color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def photo(slide, fname, x, y, w, h=None):
    path = os.path.join(BASE, fname)
    return slide.shapes.add_picture(path, x, y, w, h) if h \
           else slide.shapes.add_picture(path, x, y, w)

# ── Compound components ──────────────────────────────────────────────────────

def header(slide, title, subtitle=None, accent=GOLD):
    """Dark navy header band + accent underline."""
    R(slide, 0, 0, W, Inches(1.12), NAVY)
    R(slide, 0, Inches(1.12), W, Inches(0.055), accent)
    T(slide, title, Inches(0.5), Inches(0.15), Inches(10.5), Inches(0.7),
      28, bold=True, color=WHITE)
    if subtitle:
        T(slide, subtitle, Inches(0.5), Inches(0.8), Inches(10.5), Inches(0.36),
          12, color=RGBColor(0x88, 0xA4, 0xCC))

def card(slide, x, y, w, h, accent=BLUE, bg=WHITE):
    """White card with left accent bar."""
    R(slide, x, y, w, h, bg)
    R(slide, x, y, Inches(0.07), h, accent)
    return x, y, w, h

def card_title(slide, text, x, y, w, color=BLUE):
    T(slide, text, x + Inches(0.18), y + Inches(0.13),
      w - Inches(0.25), Inches(0.38), 12, bold=True, color=color)

def card_body(slide, lines, x, w, start_y, gap=Inches(0.3), size=10):
    cy = start_y
    for line in lines:
        T(slide, line, x + Inches(0.22), cy, w - Inches(0.32),
          Inches(0.32), size, color=SLATE)
        cy += gap
    return cy

def bullet_card(slide, x, y, w, h, title, bullets, accent=BLUE):
    """Card with title + bullet list."""
    card(slide, x, y, w, h, accent)
    card_title(slide, title, x, y, w, accent)
    R(slide, x + Inches(0.18), y + Inches(0.56),
      w - Inches(0.26), Inches(0.02), MGRAY)
    card_body(slide, [f"• {b}" for b in bullets],
              x, w, y + Inches(0.65), Inches(0.34))

def step_card(slide, num, title, body, x, y, w, h, accent=BLUE):
    """Numbered step card."""
    R(slide, x, y, w, h, WHITE)
    R(slide, x, y, w, Inches(0.055), accent)
    sz = Inches(0.46)
    OV(slide, x + Inches(0.18), y + Inches(0.16), sz, sz, accent)
    T(slide, str(num), x + Inches(0.18), y + Inches(0.2), sz, sz - Inches(0.08),
      18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(slide, title, x + Inches(0.8), y + Inches(0.19),
      w - Inches(0.95), Inches(0.38), 12, bold=True, color=NAVY)
    T(slide, body, x + Inches(0.18), y + Inches(0.73),
      w - Inches(0.3), h - Inches(0.85), 10, color=SLATE)

def rule_box(slide, x, y, w, h, text, accent=GOLD, bg=LGOLD):
    """Highlighted rule / warning / note."""
    R(slide, x, y, w, h, bg)
    R(slide, x, y, Inches(0.07), h, accent)
    T(slide, text, x + Inches(0.18), y + Inches(0.1),
      w - Inches(0.26), h - Inches(0.2), 10, color=DGRAY)

def chip(slide, label, x, y, color):
    """Small colored badge."""
    R(slide, x, y, Inches(1.38), Inches(0.3), color)
    T(slide, label, x, y + Inches(0.04), Inches(1.38), Inches(0.24),
      9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def role_chips(slide, roles):
    """Role access badges — top right corner."""
    labels = {"S": ("Superadmin", RED), "A": ("Admin", BLUE), "C": ("Cashier", GREEN)}
    rx = W - Inches(0.3) - len(roles) * Inches(1.46)
    for r in roles:
        lbl, col = labels[r]
        chip(slide, lbl, rx, Inches(0.25), col)
        rx += Inches(1.46)

def divider_slide(num, title, subtitle, accent=GOLD):
    """Full-bleed dark section divider."""
    slide = prs.slides.add_slide(blank)
    R(slide, 0, 0, W, H, NAVY)
    # Ethiopian flag strips
    R(slide, 0, 0,          Inches(0.14), H * 0.33, GREEN)
    R(slide, 0, H * 0.33,   Inches(0.14), H * 0.34, GOLD)
    R(slide, 0, H * 0.67,   Inches(0.14), H * 0.33, RED)
    # Faint huge number
    T(slide, str(num), Inches(7.0), Inches(-0.8), Inches(7.5), Inches(7.5),
      260, bold=True, color=RGBColor(0x16, 0x2A, 0x44), align=PP_ALIGN.LEFT)
    # Accent bar
    R(slide, Inches(0.5), Inches(2.0), Inches(0.08), Inches(2.5), accent)
    T(slide, f"Section {num}", Inches(0.76), Inches(2.05),
      Inches(9), Inches(0.5), 16, color=accent)
    T(slide, title, Inches(0.76), Inches(2.6),
      Inches(10), Inches(1.2), 48, bold=True, color=WHITE)
    R(slide, Inches(0.76), Inches(3.95), Inches(5), Inches(0.055), accent)
    T(slide, subtitle, Inches(0.76), Inches(4.1),
      Inches(10), Inches(0.65), 15, color=RGBColor(0x88, 0xA4, 0xCC))
    return slide

def row_table(slide, x, y, w, rows, col1_w=Inches(3.5)):
    """Alternating-row key/value table."""
    rh = Inches(0.42)
    for i, (label, value) in enumerate(rows):
        bg = LGRAY if i % 2 == 0 else WHITE
        R(slide, x, y, w, rh, bg)
        T(slide, label, x + Inches(0.15), y + Inches(0.08),
          col1_w - Inches(0.2), Inches(0.26), 10, bold=True, color=NAVY)
        T(slide, value, x + col1_w + Inches(0.1), y + Inches(0.08),
          w - col1_w - Inches(0.2), Inches(0.26), 10, color=SLATE)
        y += rh
    return y

def photo_frame(slide, fname, x, y, w, h):
    """Photo with a thin dark border."""
    R(slide, x - Inches(0.03), y - Inches(0.03),
      w + Inches(0.06), h + Inches(0.06), MGRAY)
    photo(slide, fname, x, y, w, h)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
def s_cover():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, NAVY)
    # Ethiopian flag left strip
    R(sl, 0, 0,          Inches(0.14), H*0.33, GREEN)
    R(sl, 0, H*0.33,     Inches(0.14), H*0.34, GOLD)
    R(sl, 0, H*0.67,     Inches(0.14), H*0.33, RED)
    # Right panel — screenshot
    R(sl, Inches(8.1), 0, Inches(5.23), H, RGBColor(0x09, 0x18, 0x2C))
    photo_frame(sl, "check_superadmin.png",
                Inches(8.25), Inches(0.45), Inches(4.9), Inches(5.85))
    # App icon circle
    OV(sl, Inches(0.55), Inches(0.65), Inches(1.2), Inches(1.2), GOLD)
    T(sl, "እ", Inches(0.55), Inches(0.7), Inches(1.2), Inches(1.05),
      44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # Titles
    T(sl, "እቁብ", Inches(0.55), Inches(1.9), Inches(7.2), Inches(1.5),
      82, bold=True, color=WHITE)
    T(sl, "Equb Management Platform",
      Inches(0.58), Inches(3.25), Inches(7.2), Inches(0.6), 24, color=GOLD)
    R(sl, Inches(0.58), Inches(3.95), Inches(4.5), Inches(0.055), GOLD)
    T(sl, "User & Administrator Manual",
      Inches(0.58), Inches(4.1), Inches(7.2), Inches(0.5), 18, bold=True, color=WHITE)
    T(sl, "Covers all roles: Superadmin  ·  Admin  ·  Cashier",
      Inches(0.58), Inches(4.68), Inches(7.2), Inches(0.4),
      13, color=RGBColor(0x88, 0xA4, 0xCC))
    # Stat chips (v1 style)
    chips_data = [
        ("v3.0", GREEN), ("FastAPI + PostgreSQL", BLUE),
        ("Railway.app", RED), ("30 Slides", GOLD),
    ]
    bx = Inches(0.58)
    for label, col in chips_data:
        cw = Inches(2.1) if len(label) > 8 else Inches(1.5)
        R(sl, bx, Inches(5.65), cw, Inches(1.25), RGBColor(0x0C, 0x1E, 0x34))
        R(sl, bx, Inches(5.65), cw, Inches(0.06), col)
        T(sl, label, bx, Inches(5.78), cw, Inches(0.6), 13,
          bold=True, color=col, align=PP_ALIGN.CENTER)
        bx += cw + Inches(0.12)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════════════════════
def s_toc():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "Table of Contents",
           "This manual covers all 8 functional areas of the Equb platform")

    sections = [
        ("1", "Getting Started",       "Login · Dashboard · Navigation",         GREEN),
        ("2", "Managing Members",      "Add · Edit · Spots · Lifecycle",          BLUE),
        ("3", "Cycles & Draws",        "Create cycle · Run draws · Sell pot",     GOLD),
        ("4", "Payments",              "Record · Batch · Auto-status rules",      RED),
        ("5", "Disbursements",         "Cheques · Guarantors · Vouchers",         PURPLE),
        ("6", "Reports & Analytics",   "Financial · Ledger · Export",             CYAN),
        ("7", "SMS Notifications",     "Templates · Auto-send · Log",             GREEN),
        ("8", "Settings & Admin",      "Roles · Permissions · Security · 2FA",   BLUE),
    ]
    cw = Inches(3.08); ch = Inches(2.3)
    gx = Inches(0.16); gy = Inches(0.18)
    sx = Inches(0.37); sy = Inches(1.33)
    for i, (num, title, desc, accent) in enumerate(sections):
        row, col = divmod(i, 4)
        x = sx + col*(cw+gx); y = sy + row*(ch+gy)
        card(sl, x, y, cw, ch, accent)
        # number badge
        OV(sl, x+Inches(0.2), y+Inches(0.2), Inches(0.52), Inches(0.52), accent)
        T(sl, num, x+Inches(0.2), y+Inches(0.24), Inches(0.52), Inches(0.44),
          20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        T(sl, title, x+Inches(0.88), y+Inches(0.23),
          cw-Inches(1.02), Inches(0.42), 13, bold=True, color=NAVY)
        R(sl, x+Inches(0.2), y+Inches(0.82), cw-Inches(0.4), Inches(0.02), MGRAY)
        T(sl, desc, x+Inches(0.2), y+Inches(0.94),
          cw-Inches(0.4), Inches(1.2), 10, color=SLATE)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — GETTING STARTED
# ═══════════════════════════════════════════════════════════════════════════════
def s_login():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "1.1  Logging In",
           "Authentication · Rate Limiting · Two-Factor Authentication", GREEN)
    role_chips(sl, ["S","A","C"])

    steps = [
        ("Open the app URL",
         "Navigate to the Equb platform link.\nDefault local: http://localhost:8000"),
        ("Enter your username",
         "Case-sensitive. Set by the Superadmin\nwhen your account was created."),
        ("Enter your password",
         "Min 8 characters. Cannot be all\nnumbers or a common phrase."),
        ("Click Sign In",
         "Redirected to Dashboard.\nIf 2FA is enabled, a 6-digit code\nscreen appears next."),
    ]
    sw = Inches(3.04); sh = Inches(2.0); gx = Inches(0.16); sy = Inches(1.3)
    for i, (t, d) in enumerate(steps):
        step_card(sl, i+1, t, d, Inches(0.35)+i*(sw+gx), sy, sw, sh, BLUE)

    # Right column rules
    rule_box(sl, Inches(0.35), Inches(3.5), Inches(5.8), Inches(0.95),
             "Rate Limit: After 5 failed login attempts in 5 minutes, the account is "
             "locked for 10 minutes. Contact Superadmin if you are stuck.", RED, LRED)

    rule_box(sl, Inches(0.35), Inches(4.58), Inches(5.8), Inches(0.95),
             "2FA (Optional): If enabled, after password verification you will be "
             "prompted for a 6-digit TOTP code from your authenticator app "
             "(Google Authenticator, Authy, etc.).", GOLD, LGOLD)

    bullet_card(sl, Inches(6.5), Inches(3.5), Inches(6.48), Inches(1.05),
                "Password Rules",
                ["Minimum 8 characters",
                 'Not all numbers (e.g. "12345678")',
                 'Not common phrases ("password", "equb1234")'],
                NAVY)

    rule_box(sl, Inches(6.5), Inches(4.68), Inches(6.48), Inches(0.85),
             "Forgot your password? Ask your Superadmin — they can reset any "
             "user password from Settings → Users.", BLUE, LBLUE)

    rule_box(sl, Inches(0.35), Inches(5.7), Inches(12.63), Inches(0.88),
             "First login: Your Superadmin will give you a username and temporary password. "
             "Change your password immediately via Settings → Change Password.", GREEN, LGREEN)


def s_dashboard():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "1.2  Dashboard Overview",
           "Real-time snapshot of the active cycle — visible to all roles", GREEN)
    role_chips(sl, ["S","A","C"])

    photo_frame(sl, "check_superadmin.png",
                Inches(5.35), Inches(1.28), Inches(7.65), Inches(5.95))

    callouts = [
        (GREEN, Inches(1.65), "Sidebar Navigation",
         "Links: Dashboard · Members\nDraws · Payments · Reports\nNotifications · Settings"),
        (GOLD,  Inches(2.9),  "Active Cycle Info",
         "Shows which cycle is selected\nand the current draw phase\n(Collection or Active)"),
        (BLUE,  Inches(4.15), "Next Draw Card",
         "Week number, scheduled\ndraw date, and net pot\namount for that week"),
        (RED,   Inches(5.4),  "Association Fund",
         "Running ማህበር fund balance\nbuilt from weekly member\ndeductions"),
    ]
    for accent, cy, title, desc in callouts:
        OV(sl, Inches(0.32), cy-Inches(0.18), Inches(0.35), Inches(0.35), accent)
        R(sl, Inches(0.67), cy-Inches(0.01), Inches(0.38), Inches(0.03), MGRAY)
        T(sl, title, Inches(1.12), cy-Inches(0.22),
          Inches(4.0), Inches(0.38), 12, bold=True, color=NAVY)
        T(sl, desc, Inches(1.12), cy+Inches(0.2),
          Inches(4.0), Inches(0.78), 10, color=SLATE)
        R(sl, Inches(0.32), cy+Inches(1.0), Inches(4.85), Inches(0.02), MGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 2 — MEMBERS
# ═══════════════════════════════════════════════════════════════════════════════
def s_member_concepts():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "2.1  Member Concepts",
           "Spots · Full Share vs Half Share · Weekly Contribution", BLUE)
    role_chips(sl, ["S","A"])

    bullet_card(sl, Inches(0.35), Inches(1.3), Inches(6.0), Inches(2.0),
                "What is a Spot?",
                ["Each spot is a numbered slot in the Equb (e.g. #1 to #113)",
                 "One spot = one draw entry per cycle",
                 "Member spots: #1–113  |  Association (ማህበር) spots: #114–118",
                 "Association spots generate profit for the group fund",
                 "A member can hold more than one spot"], GREEN)

    # Full vs Half side-by-side
    for col, (label, acc, lines) in enumerate([
        ("Full Share  —  21,000 ETB / week", GREEN, [
            "One member occupies the spot exclusively",
            "Receives the full pot when drawn",
            "Cannot be split with another member",
        ]),
        ("Half Share  —  10,500 ETB / week", GOLD, [
            "Two members share ONE spot",
            "Each pays 10,500 ETB per week",
            "When drawn, pot is split between both",
            "Maximum 2 members per half-share spot",
        ]),
    ]):
        x = Inches(0.35) + col*Inches(3.16)
        card(sl, x, Inches(3.5), Inches(3.0), Inches(2.62), acc)
        T(sl, label, x+Inches(0.18), Inches(3.63),
          Inches(2.76), Inches(0.4), 11, bold=True, color=acc)
        R(sl, x+Inches(0.18), Inches(4.08), Inches(2.68), Inches(0.02), MGRAY)
        cy = Inches(4.18)
        for ln in lines:
            T(sl, f"• {ln}", x+Inches(0.22), cy, Inches(2.6), Inches(0.32), 10, color=SLATE)
            cy += Inches(0.36)

    # Settings reference table — right
    R(sl, Inches(6.6), Inches(1.3), Inches(6.38), Inches(0.4), NAVY)
    T(sl, "Default Financial Settings", Inches(6.75), Inches(1.38),
      Inches(6.1), Inches(0.28), 12, bold=True, color=WHITE)
    row_table(sl, Inches(6.6), Inches(1.7), Inches(6.38), [
        ("full_spot_amount",      "21,000 ETB / week"),
        ("half_spot_amount",      "10,500 ETB / week"),
        ("association_deduction", "1,000 ETB / person / week"),
        ("full_spot_voucher",     "80 ETB / card"),
        ("half_spot_voucher",     "40 ETB / card"),
        ("group_week_interval",   "Every 4th week is a buyer week"),
        ("total_member_spots",    "113 regular member spots"),
        ("total_assoc_spots",     "5 association spots"),
    ], Inches(3.2))
    rule_box(sl, Inches(6.6), Inches(5.25), Inches(6.38), Inches(0.82),
             "Each cycle takes a snapshot of these values at creation. "
             "Settings changes only affect NEW cycles unless you use 'Apply Settings to Cycle'.",
             GOLD, LGOLD)


def s_add_member():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "2.2  Adding & Editing Members",
           "Members page → + Add Member button", BLUE)
    role_chips(sl, ["S","A"])

    steps = [
        ("Go to Members page",
         "Click Members in the sidebar.\nClick the + Add Member button\nat the top right."),
        ("Enter member details",
         "Name (required)\nPhone number (optional, for SMS)\nNotes (optional)"),
        ("Assign a spot",
         "Select spot number from the\navailable list.\nChoose share type: Full or Half."),
        ("Save the member",
         "Click Save Member.\nMember appears in the list\nimmediately."),
        ("Add more spots (optional)",
         "Open the member record.\nClick Add Spot to assign\nan additional spot."),
        ("Bulk import",
         "Click Import to upload a CSV\nor Excel file.\nColumns: Name, Phone, Spot, Share, Notes."),
    ]
    sw = Inches(4.08); sh = Inches(1.93); gx = Inches(0.16); gy = Inches(0.16)
    sx = Inches(0.35); sy = Inches(1.3)
    for i, (t, d) in enumerate(steps):
        row, col = divmod(i, 3)
        step_card(sl, i+1, t, d,
                  sx+col*(sw+gx), sy+row*(sh+gy), sw, sh, BLUE)

    rule_box(sl, Inches(0.35), Inches(5.65), Inches(12.63), Inches(0.88),
             "Important: Spot assignments cannot be added or edited once the cycle moves "
             "to the Draw phase. Complete all member registrations during the Collection phase.",
             RED, LRED)


def s_member_lifecycle():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "2.3  Member Lifecycle & Status",
           "Active → Received → Left  |  Mid-cycle exit  |  Spot transfer", BLUE)
    role_chips(sl, ["S","A"])

    # Status boxes
    statuses = [
        ("ACTIVE",    GREEN,  "Member is currently\nparticipating and\npaying weekly"),
        ("RECEIVED",  GOLD,   "Member has won the\npot and collected\nthe cheque"),
        ("LEFT",      RED,    "Member has departed\nor stopped paying\nduring the cycle"),
    ]
    bw = Inches(3.5); bh = Inches(1.6)
    sx = Inches(0.5); sy = Inches(1.3)
    for i, (lbl, acc, desc) in enumerate(statuses):
        bx = sx + i*(bw+Inches(0.58))
        R(sl, bx, sy, bw, bh, acc)
        T(sl, lbl, bx, sy+Inches(0.2), bw, Inches(0.44),
          20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        T(sl, desc, bx, sy+Inches(0.72), bw, Inches(0.76),
          11, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 2:
            R(sl, bx+bw+Inches(0.14), sy+Inches(0.76), Inches(0.3), Inches(0.04), MGRAY)
            # arrow tip
            tip = sl.shapes.add_shape(13, bx+bw+Inches(0.44), sy+Inches(0.68),
                                       Inches(0.16), Inches(0.2))
            tip.fill.solid(); tip.fill.fore_color.rgb = MGRAY
            tip.line.fill.background()

    bullet_card(sl, Inches(0.35), Inches(3.2), Inches(6.1), Inches(2.4),
                "Mid-Cycle Exit (Member Leaves)",
                ["Go to Members → open member → click Exit Member",
                 "Enter the exit week number and reason (left / stopped paying)",
                 "All pending/late payments AFTER exit week → auto-marked missed",
                 "A financial summary is shown: paid, owed, whether they won",
                 "Soft delete — all records preserved for audit"], RED)

    bullet_card(sl, Inches(6.6), Inches(3.2), Inches(6.38), Inches(2.4),
                "Spot Transfer (Mid-Cycle)",
                ["Go to Members → click Spot Transfer",
                 "Select spot, current member (outgoing), new member (incoming)",
                 "Specify effective week (first week new member pays)",
                 "Past payments stay with original member",
                 "Future pending payments reassigned to new member"], BLUE)

    rule_box(sl, Inches(0.35), Inches(5.78), Inches(12.63), Inches(0.78),
             "Members marked Left are never permanently deleted. "
             "All payment, draw, and guarantor records are preserved for audit trail.",
             GOLD, LGOLD)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 3 — CYCLES & DRAWS
# ═══════════════════════════════════════════════════════════════════════════════
def s_create_cycle():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "3.1  Creating & Managing a Cycle",
           "Draws page → Cycles tab → + New Cycle", GOLD)
    role_chips(sl, ["S","A"])

    steps = [
        ("Go to Draws → Cycles",
         'Click Draws in sidebar.\nSelect the Cycles tab.\nClick "+ New Cycle".'),
        ("Fill cycle details",
         "Name (e.g. Equb 2025)\nStart date\nNotes (optional)"),
        ("Override settings (optional)",
         "Per-cycle override of: spot amounts,\nassoc. deduction, voucher rates,\ntotal member spots."),
        ("Create the cycle",
         "Click Create. All weeks are\ngenerated automatically.\nStatus: Collection phase."),
        ("Register all members",
         "Add members & assign spots\nduring Collection phase.\nNo draws happen yet."),
        ("Start Draws phase",
         'Click "Start Draws".\nSpecify first week number and\ntotal association spots.'),
    ]
    sw = Inches(4.08); sh = Inches(1.93); gx = Inches(0.16); gy = Inches(0.16)
    sx = Inches(0.35); sy = Inches(1.3)
    for i, (t, d) in enumerate(steps):
        row, col = divmod(i, 3)
        step_card(sl, i+1, t, d,
                  sx+col*(sw+gx), sy+row*(sh+gy), sw, sh, GREEN)

    rule_box(sl, Inches(0.35), Inches(5.65), Inches(12.63), Inches(0.88),
             "Only one cycle can be active at a time. Creating a new cycle automatically "
             "closes the previous active cycle. Completed cycles are locked and preserved for audit.",
             GOLD, LGOLD)


def s_draws():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "3.2  Running Weekly Draws & Selling the Pot",
           "Draws page → Weeks tab", GOLD)
    role_chips(sl, ["S","A"])

    # Draw steps — left
    draw_steps = [
        ("Go to Draws → Weeks",
         "Find the current week row\nin the weeks list."),
        ("Click Draw Winner",
         'Click the "Draw" button.\nSelect or confirm the winning spot.'),
        ("Eligibility check",
         "Winner must have no pending,\nlate, or missed payments."),
        ("Draw confirmed",
         "Spot → Received. Week → Drawn.\nSMS confirmation sent to winner."),
    ]
    sy = Inches(1.3)
    for i, (t, d) in enumerate(draw_steps):
        step_card(sl, i+1, t, d, Inches(0.35), sy+i*Inches(1.42), Inches(5.8), Inches(1.3), BLUE)

    # Right column
    card(sl, Inches(6.55), Inches(1.3), Inches(6.43), Inches(2.8), RED)
    T(sl, "Pot On Hold", Inches(6.73), Inches(1.43),
      Inches(6.1), Inches(0.38), 13, bold=True, color=RED)
    R(sl, Inches(6.73), Inches(1.86), Inches(6.1), Inches(0.02), MGRAY)
    for j, line in enumerate([
        "• Winner has unpaid weeks → pot placed ON HOLD",
        "• SMS sent: 'Pay outstanding weeks to receive pot'",
        "• Member clears all outstanding payments",
        "• Admin releases the pot manually",
        "• Cheque issued the following week",
    ]):
        T(sl, line, Inches(6.73), Inches(1.98)+j*Inches(0.4),
          Inches(6.1), Inches(0.35), 10, color=SLATE)

    bullet_card(sl, Inches(6.55), Inches(4.3), Inches(6.43), Inches(1.85),
                "Pot Sales — 3 Options",
                ["Group Week (every 4th week): any member can buy the pot",
                 "Member Sale: winning member sells their pot to another member",
                 "Association Spot Sale: ማህበር spot sold; profit → assoc fund",
                 "Go to Draws → select week → Sell Pot"], GOLD)

    rule_box(sl, Inches(0.35), Inches(6.88), Inches(5.8), Inches(0.68),
             "Batch Draw: Draw multiple weeks at once from Draws → Batch Draw. "
             "Each week is validated independently.", BLUE, LBLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 4 — PAYMENTS
# ═══════════════════════════════════════════════════════════════════════════════
def s_payments():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "4.1  Recording Payments",
           "Payments page → Select week → Record payment", RED)
    role_chips(sl, ["S","A","C"])

    bullet_card(sl, Inches(0.35), Inches(1.3), Inches(5.9), Inches(3.1),
                "Record a Single Payment",
                ["Go to Payments → select the week from dropdown",
                 "Find the member in the payment list",
                 "Click Record Payment next to their name",
                 "Fill in: Status, Date, Method, Reference, Notes",
                 "Click Save — SMS confirmation sent automatically"], BLUE)

    bullet_card(sl, Inches(6.6), Inches(1.3), Inches(6.38), Inches(3.1),
                "Record a Batch Payment (Multiple Weeks)",
                ["Go to Payments → Batch Record",
                 "Select the member from the dropdown",
                 "Tick the week checkboxes you are recording for",
                 "Enter payment date, method, and reference",
                 "Click Save — all selected weeks are marked paid",
                 "One SMS sent summarising all weeks covered"], GREEN)

    # Fields reference
    R(sl, Inches(0.35), Inches(4.55), Inches(12.63), Inches(0.38), NAVY)
    T(sl, "Payment Fields Reference", Inches(0.5), Inches(4.6),
      Inches(12.0), Inches(0.28), 12, bold=True, color=WHITE)
    row_table(sl, Inches(0.35), Inches(4.93), Inches(12.63), [
        ("Status",          "pending  ·  paid  ·  late  ·  missed"),
        ("Payment date",    "Date cash / cheque was physically received"),
        ("Payment method",  "cash  ·  bank_transfer  ·  cheque"),
        ("Reference",       "Cheque number or bank transaction reference"),
        ("Penalty amount",  "ETB surcharge for late payment (optional)"),
        ("Notes",           "Any additional remarks"),
    ], Inches(3.2))


def s_payment_rules():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "4.2  Payment Status Rules",
           "Automatic transitions · Scheduled jobs · Manual override", RED)
    role_chips(sl, ["S","A","C"])

    # Status flow
    statuses = [
        ("PENDING", BLUE,   "Payment created\nbut not yet\nrecorded"),
        ("PAID",    GREEN,  "Payment confirmed\nby cashier or\nadmin"),
        ("LATE",    GOLD,   "Draw date passed\n(within 3-day\ngrace period)"),
        ("MISSED",  RED,    "3+ days past draw\ndate without\npayment"),
    ]
    bw = Inches(2.9); bh = Inches(1.65)
    sx = Inches(0.38); sy = Inches(1.3)
    for i, (lbl, acc, desc) in enumerate(statuses):
        bx = sx + i*(bw+Inches(0.22))
        R(sl, bx, sy, bw, bh, acc)
        T(sl, lbl, bx, sy+Inches(0.18), bw, Inches(0.42),
          20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        T(sl, desc, bx, sy+Inches(0.66), bw, Inches(0.86),
          11, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            R(sl, bx+bw+Inches(0.04), sy+Inches(0.78),
              Inches(0.14), Inches(0.04), MGRAY)

    bullet_card(sl, Inches(0.35), Inches(3.18), Inches(6.1), Inches(2.35),
                "Nightly Automated Job (midnight EAT)",
                ["pending → late: draw date passed, within 3-day grace period",
                 "late → missed: 3 or more days past draw date",
                 "pending → missed: pending and more than 3 days past draw",
                 "SMS alert sent to each newly missed member"], NAVY)

    bullet_card(sl, Inches(6.6), Inches(3.18), Inches(6.38), Inches(2.35),
                "Daily Reminder Job (9 PM EAT)",
                ["Checks all unpaid members with draws within 48 hours",
                 "Sends payment_reminder SMS to each member",
                 "Encourages payment to maintain draw eligibility",
                 "Can also be triggered manually from Payments page"], CYAN)

    rule_box(sl, Inches(0.35), Inches(5.68), Inches(12.63), Inches(0.78),
             "Manual override: Any admin or cashier can change a payment status manually "
             "at any time. Go to Payments → select the week → click the status badge.",
             GOLD, LGOLD)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 5 — DISBURSEMENTS
# ═══════════════════════════════════════════════════════════════════════════════
def s_disbursements():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "5.1  Issuing a Disbursement Cheque",
           "Disbursements page → New Disbursement", PURPLE)
    role_chips(sl, ["S","A","C"])

    steps = [
        ("Verify winner is paid",
         "Member must have no pending,\nlate, or missed payments before\nissuing the cheque."),
        ("Go to Disbursements",
         "Click Disbursements in sidebar.\nClick New Disbursement."),
        ("Select the week",
         "Choose the draw week.\nThe gross pot amount\nauto-fills."),
        ("Get voucher info",
         'Click "Get Voucher Info" to see\nthe full deduction breakdown:\nassoc + service + voucher.'),
        ("Add 3 guarantors",
         "Select 3 different members.\nAll 3 must be different\npeople (system validates)."),
        ("Enter cheque details",
         "Cheque number and date.\nAdjust deductions if needed.\nClick Save."),
    ]
    sw = Inches(4.08); sh = Inches(1.93); gx = Inches(0.16); gy = Inches(0.16)
    sx = Inches(0.35); sy = Inches(1.3)
    for i, (t, d) in enumerate(steps):
        row, col = divmod(i, 3)
        step_card(sl, i+1, t, d,
                  sx+col*(sw+gx), sy+row*(sh+gy), sw, sh, PURPLE)

    rule_box(sl, Inches(0.35), Inches(5.65), Inches(12.63), Inches(0.88),
             "3 Guarantors Required — all three must be different members. "
             "The system rejects if any two guarantors are the same person. "
             "Guarantor records are preserved permanently in the audit trail.",
             RED, LRED)


def s_voucher():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "5.2  Net Pot Calculation & Voucher Deductions",
           "How the winner's cheque amount is calculated", PURPLE)
    role_chips(sl, ["S","A"])

    # Formula bar
    R(sl, Inches(0.35), Inches(1.3), Inches(12.63), Inches(0.65), NAVY)
    T(sl, "Net Pot  =  Gross Pot  −  Association Deduction  −  Service Fee  −  Voucher Deduction",
      Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.44),
      20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Worked example table
    R(sl, Inches(0.35), Inches(2.05), Inches(12.63), Inches(0.38), DGRAY)
    T(sl, "Worked Example — Full-spot winner, 113 members, 120-week cycle",
      Inches(0.5), Inches(2.1), Inches(12.2), Inches(0.28),
      11, bold=True, color=WHITE)
    example_rows = [
        ("Gross Pot",              "113 members × 21,000 ETB/week",          "2,373,000 ETB"),
        ("− Association Deduction","113 members × 1,000 ETB (assoc. rate)",  "  − 113,000 ETB"),
        ("− Service Fee",          "1 worker-week × 21,000 ETB",             "   − 21,000 ETB"),
        ("− Voucher Deduction",    "1 winner × 80 ETB/card × 120 weeks",     "    − 9,600 ETB"),
    ]
    ry = Inches(2.43)
    for i, (lbl, formula, amount) in enumerate(example_rows):
        bg = LGRAY if i%2==0 else WHITE
        R(sl, Inches(0.35), ry, Inches(12.63), Inches(0.48), bg)
        T(sl, lbl,     Inches(0.5), ry+Inches(0.09), Inches(3.8), Inches(0.3), 10, bold=True, color=NAVY)
        T(sl, formula, Inches(4.35), ry+Inches(0.09), Inches(6.0), Inches(0.3), 10, color=SLATE)
        T(sl, amount,  Inches(10.4), ry+Inches(0.09), Inches(2.3), Inches(0.3), 10,
          color=RED, bold=True, align=PP_ALIGN.RIGHT)
        ry += Inches(0.48)
    # Net total
    R(sl, Inches(0.35), ry, Inches(12.63), Inches(0.52), NAVY)
    T(sl, "NET CHEQUE AMOUNT", Inches(0.5), ry+Inches(0.1), Inches(5.0), Inches(0.34),
      13, bold=True, color=WHITE)
    T(sl, "2,229,400 ETB", Inches(10.4), ry+Inches(0.1), Inches(2.3), Inches(0.34),
      14, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

    bullet_card(sl, Inches(0.35), Inches(4.9), Inches(6.1), Inches(1.82),
                "Voucher Returns (Physical Cards)",
                ["Go to Disbursements → Voucher Returns → Add Return",
                 "Enter full_count and half_count of cards returned by vendor",
                 "Mark vendor_paid when vendor is reimbursed",
                 "One record per week — cannot be duplicated"], CYAN)

    bullet_card(sl, Inches(6.6), Inches(4.9), Inches(6.38), Inches(1.82),
                "Void a Disbursement",
                ["Open the disbursement record → click Void",
                 "Enter a reason (minimum 5 characters)",
                 "Status changes to Voided — soft delete only",
                 "All original data preserved for audit trail",
                 "Records who voided it and when"], RED)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 6 — REPORTS
# ═══════════════════════════════════════════════════════════════════════════════
def s_reports():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "6.1  Reports & Analytics",
           "Reports page — Admin / Superadmin access required", CYAN)
    role_chips(sl, ["S","A"])

    photo_frame(sl, "pg_reports.png",
                Inches(6.6), Inches(1.28), Inches(6.4), Inches(5.95))

    reports = [
        (GOLD,   "Financial Summary",
         "Tabs: Summary · Balance Sheet · Vouchers · General Ledger. "
         "Total collected, assoc fund, service fee, collection rate."),
        (GREEN,  "Payment Report",
         "Per-member: weeks paid/missed, total amounts. Full payment "
         "history with date, method, and collector name."),
        (BLUE,   "Draw Report",
         "All completed draws — winner, buyer (if sold), pot amounts, "
         "draw date, and pot-sale details."),
        (RED,    "Disbursement Report",
         "All cheques: gross/net/deductions, cheque #, guarantors, "
         "status (issued/collected/voided), void reasons."),
        (PURPLE, "Association Fund Report",
         "Inflows (deductions + ማህበር profits) minus expenses. "
         "Running balance and end-of-cycle distribution."),
    ]
    y = Inches(1.38)
    for acc, title, desc in reports:
        OV(sl, Inches(0.32), y+Inches(0.04), Inches(0.34), Inches(0.34), acc)
        T(sl, title, Inches(0.8), y, Inches(5.6), Inches(0.36), 12, bold=True, color=NAVY)
        T(sl, desc, Inches(0.8), y+Inches(0.37), Inches(5.6), Inches(0.65), 10, color=SLATE)
        R(sl, Inches(0.32), y+Inches(1.08), Inches(6.1), Inches(0.02), MGRAY)
        y += Inches(1.14)

    rule_box(sl, Inches(0.32), Inches(6.32), Inches(6.1), Inches(0.72),
             "Export: All reports download as CSV or formatted Excel. "
             "Click Print Report to open a print-ready PDF summary.", GOLD, LGOLD)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 7 — NOTIFICATIONS
# ═══════════════════════════════════════════════════════════════════════════════
def s_notifications():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "7.1  SMS Notifications Setup",
           "Notifications page → Settings tab  |  Africa's Talking API", GREEN)
    role_chips(sl, ["S","A"])

    photo_frame(sl, "pg_notifications.png",
                Inches(6.6), Inches(1.28), Inches(6.4), Inches(5.95))

    setup_steps = [
        ("Open Notifications → Settings",
         "Click Notifications in sidebar.\nClick the Settings tab."),
        ("Enter Africa's Talking credentials",
         "API Key, Username (sandbox or live).\nSender ID (your SMS display name)."),
        ("Enable SMS",
         'Toggle "Enable SMS" to ON.\nLeave OFF for testing (mock mode\n— logs without sending).'),
        ("Save & test",
         "Click Save Settings.\nUse the test send button to\nverify delivery."),
    ]
    sy = Inches(1.32)
    for i, (t, d) in enumerate(setup_steps):
        step_card(sl, i+1, t, d,
                  Inches(0.32), sy+i*Inches(1.42), Inches(6.1), Inches(1.3), GREEN)

    rule_box(sl, Inches(0.32), Inches(7.0), Inches(6.1), Inches(0.58),
             "Mock Mode: SMS disabled → notifications logged as 'mock'. "
             "See exactly what would be sent without using credits.", GOLD, LGOLD)


def s_sms_templates():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "7.2  SMS Templates",
           "6 built-in templates — editable, individually toggleable", GREEN)
    role_chips(sl, ["S","A"])

    templates = [
        (GREEN,  "payment_confirmed",
         "Sent when cashier marks a payment paid.",
         "{member_name}  {amount}  {week_number}  {draw_date}  {payment_method}"),
        (GOLD,   "payment_reminder",
         "Daily reminder — unpaid members with draw within 48 h.",
         "{member_name}  {amount}  {week_number}  {draw_date}"),
        (RED,    "missed_payment",
         "Sent when payment auto-transitions to missed.",
         "{member_name}  {unpaid_count}  {unpaid_amount}  {weeks_list}"),
        (BLUE,   "draw_winner",
         "Congratulations SMS to the draw winner.",
         "{member_name}  {week_number}  {net_pot}"),
        (PURPLE, "pot_on_hold",
         "Winner drawn but has outstanding unpaid weeks.",
         "{member_name}  {week_number}  {net_pot}  {unpaid_count}  {unpaid_amount}"),
        (CYAN,   "pot_sold",
         "Sent when a pot is sold instead of directly disbursed.",
         "{member_name}  {week_number}  {seller_fee}  {buyer_receives}"),
    ]
    cw = Inches(4.08); ch = Inches(1.78); gx = Inches(0.16); gy = Inches(0.16)
    sx = Inches(0.35); sy = Inches(1.32)
    for i, (acc, key, trigger, variables) in enumerate(templates):
        row, col = divmod(i, 3)
        x = sx+col*(cw+gx); y = sy+row*(ch+gy)
        card(sl, x, y, cw, ch, acc)
        T(sl, key, x+Inches(0.18), y+Inches(0.1), cw-Inches(0.25), Inches(0.35),
          11, bold=True, color=acc)
        T(sl, trigger, x+Inches(0.18), y+Inches(0.5), cw-Inches(0.25), Inches(0.42),
          10, color=NAVY)
        R(sl, x+Inches(0.18), y+Inches(0.98), cw-Inches(0.26), Inches(0.02), MGRAY)
        T(sl, "Variables:  "+variables, x+Inches(0.18), y+Inches(1.06),
          cw-Inches(0.25), Inches(0.58), 9, color=SLATE, italic=True)

    rule_box(sl, Inches(0.35), Inches(5.67), Inches(12.63), Inches(0.88),
             "Edit any template in Notifications → Templates tab. "
             "Toggle a template off to stop that specific SMS type from being sent.",
             GOLD, LGOLD)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION 8 — SETTINGS & ADMIN
# ═══════════════════════════════════════════════════════════════════════════════
def s_financial_settings():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "8.1  Financial Settings",
           "Settings page → Financial tab  (Superadmin only)", BLUE)
    role_chips(sl, ["S"])

    photo_frame(sl, "pg_settings.png",
                Inches(6.6), Inches(1.28), Inches(6.4), Inches(5.95))

    R(sl, Inches(0.35), Inches(1.3), Inches(6.1), Inches(0.38), NAVY)
    T(sl, "Configurable Fields", Inches(0.5), Inches(1.35),
      Inches(5.8), Inches(0.28), 11, bold=True, color=WHITE)
    row_table(sl, Inches(0.35), Inches(1.68), Inches(6.1), [
        ("full_spot_amount",      "21,000 ETB — full-share weekly payment"),
        ("half_spot_amount",      "10,500 ETB — half-share weekly payment"),
        ("association_deduction", "1,000 ETB — per person per week → ማህበር fund"),
        ("full_spot_voucher",     "80 ETB — vendor voucher per full-spot card"),
        ("half_spot_voucher",     "40 ETB — vendor voucher per half-spot card"),
        ("group_week_interval",   "4 — every 4th week is a buyer/group week"),
        ("total_member_spots",    "113 — regular member spots per cycle"),
        ("group_name",            "Displayed in app header and all reports"),
        ("group_tagline",         "Subtitle shown under the group name"),
        ("logo_url",              "HTTPS URL for the group logo image"),
    ], Inches(2.85))

    rule_box(sl, Inches(0.35), Inches(6.08), Inches(6.1), Inches(0.98),
             "Changes affect NEW cycles only. Each cycle takes a snapshot at creation. "
             "To push changes into an active cycle, use Members → Apply Settings to Cycle.",
             RED, LRED)


def s_roles():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "8.2  User Roles & Permission Matrix",
           "Settings page → Users tab  |  Who can do what", BLUE)
    role_chips(sl, ["S"])

    roles_data = [
        ("Superadmin", RED, [
            "Full access to every feature",
            "Manage all user accounts and roles",
            "Configure financial settings & branding",
            "Reset the entire system",
            "Cannot create a second Superadmin",
        ]),
        ("Admin", BLUE, [
            "Manage members, draws, payments",
            "Issue and void disbursements",
            "View reports and export data",
            "Configure SMS notifications",
            "Manage Cashier accounts",
            "Permissions customisable by Superadmin",
        ]),
        ("Cashier", GREEN, [
            "Record and update payments",
            "Record disbursements",
            "View member list (read-only)",
            "Cannot manage members or run draws",
            "Cannot access reports or settings",
            "Permissions customisable by Superadmin",
        ]),
    ]
    rw = Inches(4.02); rh = Inches(3.55); gx = Inches(0.18)
    rx = Inches(0.35); ry = Inches(1.3)
    for i, (role, acc, perms) in enumerate(roles_data):
        x = rx + i*(rw+gx)
        R(sl, x, ry, rw, Inches(0.58), acc)
        T(sl, role, x, ry+Inches(0.1), rw, Inches(0.42),
          18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        card(sl, x, ry+Inches(0.58), rw, rh-Inches(0.58), acc)
        cy = ry+Inches(0.78)
        for p in perms:
            T(sl, f"• {p}", x+Inches(0.2), cy, rw-Inches(0.3), Inches(0.36), 10, color=SLATE)
            cy += Inches(0.42)

    # Permission matrix
    R(sl, Inches(0.35), Inches(5.0), Inches(12.63), Inches(0.34), NAVY)
    T(sl, "Permission Matrix  —  configurable by Superadmin in Settings → Permissions",
      Inches(0.5), Inches(5.05), Inches(12.2), Inches(0.24),
      11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    perms_matrix = [
        ("manage_members","✓","✓","—"),("run_draws","✓","✓","—"),
        ("disbursements","✓","✓","✓"),("view_reports","✓","✓","—"),
        ("manage_users","✓","✓","—"),("notifications","✓","✓","—"),
    ]
    col_ws = [Inches(4.5), Inches(2.7), Inches(2.7), Inches(2.73)]
    headers = ["Permission","Superadmin","Admin","Cashier"]
    hx = Inches(0.35)
    R(sl, hx, Inches(5.34), Inches(12.63), Inches(0.3), DGRAY)
    for hdr, cw in zip(headers, col_ws):
        T(sl, hdr, hx+Inches(0.08), Inches(5.38), cw-Inches(0.12), Inches(0.22),
          9, bold=True, color=GOLD)
        hx += cw
    ry2 = Inches(5.64)
    for ri, (perm, s, a, c) in enumerate(perms_matrix):
        mx = Inches(0.35)
        bg = LGRAY if ri%2==0 else WHITE
        for ci, (val, cw) in enumerate(zip([perm,s,a,c], col_ws)):
            R(sl, mx, ry2, cw, Inches(0.28), bg)
            col = NAVY if ci==0 else (GREEN if val=="✓" else RED)
            T(sl, val, mx+Inches(0.1), ry2+Inches(0.04),
              cw-Inches(0.15), Inches(0.22), 9,
              bold=(val in ("✓","—")), color=col)
            mx += cw
        ry2 += Inches(0.28)


def s_security():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "8.3  Security Features",
           "2FA · Audit Log · Sessions · Password Policy", BLUE)
    role_chips(sl, ["S","A","C"])

    bullet_card(sl, Inches(0.35), Inches(1.3), Inches(6.1), Inches(2.9),
                "Two-Factor Authentication (2FA) Setup",
                ["Go to Settings → Your Account → Enable 2FA",
                 "A QR code and secret key are displayed",
                 "Scan the QR code with Google Authenticator or Authy",
                 "Enter the 6-digit TOTP code to confirm and activate",
                 "Every future login requires password + TOTP code",
                 "To disable: enter current TOTP code in Settings → Disable 2FA"], BLUE)

    bullet_card(sl, Inches(6.6), Inches(1.3), Inches(6.38), Inches(2.9),
                "Audit Trail (Immutable Log)",
                ["Every financial change is recorded automatically",
                 "Logged: user, timestamp, action (create / update / delete / void)",
                 "Before and after values stored as JSON snapshots",
                 "Cannot be altered or deleted — permanent record",
                 "Covers: payments, draws, disbursements, member exits, voids",
                 "Viewable by Superadmin only"], PURPLE)

    sec_items = [
        (RED,    "Rate Limiting",     "5 failed logins → 10-minute account lockout"),
        (GOLD,   "Password Policy",   "Min 8 chars · no all-numbers · no common phrases"),
        (GREEN,  "Session Timeout",   "1-hour server-side session; re-login required"),
        (BLUE,   "CSRF Protection",   "All forms include CSRF tokens against forgery"),
        (PURPLE, "CSP Headers",       "Content Security Policy blocks XSS attacks"),
        (CYAN,   "HTTPS Enforced",    "All traffic forced to HTTPS in production"),
    ]
    cw = Inches(4.08); ch = Inches(0.98); gx = Inches(0.16); gy = Inches(0.12)
    sx = Inches(0.35); sy = Inches(4.38)
    for i, (acc, title, desc) in enumerate(sec_items):
        row, col = divmod(i, 3)
        x = sx+col*(cw+gx); y = sy+row*(ch+gy)
        card(sl, x, y, cw, ch, acc)
        T(sl, title, x+Inches(0.18), y+Inches(0.1), cw-Inches(0.25), Inches(0.34),
          11, bold=True, color=acc)
        T(sl, desc, x+Inches(0.18), y+Inches(0.5), cw-Inches(0.25), Inches(0.38),
          10, color=SLATE)


# ═══════════════════════════════════════════════════════════════════════════════
# GLOSSARY
# ═══════════════════════════════════════════════════════════════════════════════
def s_glossary():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, LGRAY)
    header(sl, "Glossary", "Key terms used throughout this manual")

    terms = [
        ("Equb (እቁብ)",          "Ethiopian rotating savings & credit association (ROSCA)"),
        ("Cycle",                "One complete round of the Equb — all spots drawn once"),
        ("Spot",                 "A numbered entry slot (#1–113 member, #114–118 assoc.)"),
        ("Full Share",           "One member holds a spot exclusively; pays 21,000 ETB/week"),
        ("Half Share",           "Two members share one spot; each pays 10,500 ETB/week"),
        ("Gross Pot",            "Total weekly contributions from all members"),
        ("Net Pot",              "Gross − Association deduction − Service fee − Voucher"),
        ("Association (ማህበር)",   "Community fund from weekly deductions"),
        ("Voucher",              "Physical vendor discount card; reimbursed by assoc. fund"),
        ("Group/Buyer Week",     "Every 4th week — any member can buy the pot"),
        ("Guarantor",            "Member vouching for the draw winner (3 required)"),
        ("Disbursement",         "Physical cheque payment issued to the draw winner"),
        ("Collection Phase",     "Early cycle stage — members register, no draws yet"),
        ("Draw Phase",           "Active stage — weekly draws are held"),
        ("Pot On Hold",          "Winner drawn but cannot receive pot due to unpaid weeks"),
        ("Void",                 "Cancellation of a disbursement cheque (soft delete)"),
        ("TOTP / 2FA",           "Time-based one-time password for two-factor authentication"),
    ]
    mid = len(terms)//2 + 1
    for ci, term_list in enumerate([terms[:mid], terms[mid:]]):
        x = Inches(0.35) + ci*Inches(6.6)
        ry = Inches(1.3)
        for i, (term, defn) in enumerate(term_list):
            bg = LGRAY if i%2==0 else WHITE
            R(sl, x, ry, Inches(6.35), Inches(0.43), bg)
            T(sl, term, x+Inches(0.15), ry+Inches(0.07),
              Inches(2.2), Inches(0.28), 10, bold=True, color=NAVY)
            T(sl, defn, x+Inches(2.4), ry+Inches(0.07),
              Inches(3.8), Inches(0.28), 10, color=SLATE)
            ry += Inches(0.43)


# ═══════════════════════════════════════════════════════════════════════════════
# BACK COVER
# ═══════════════════════════════════════════════════════════════════════════════
def s_back_cover():
    sl = prs.slides.add_slide(blank)
    R(sl, 0, 0, W, H, NAVY)
    R(sl, 0, 0,          Inches(0.14), H*0.33, GREEN)
    R(sl, 0, H*0.33,     Inches(0.14), H*0.34, GOLD)
    R(sl, 0, H*0.67,     Inches(0.14), H*0.33, RED)
    xr = W - Inches(0.14)
    R(sl, xr, 0,         Inches(0.14), H*0.33, GREEN)
    R(sl, xr, H*0.33,    Inches(0.14), H*0.34, GOLD)
    R(sl, xr, H*0.67,    Inches(0.14), H*0.33, RED)

    T(sl, "Equb Management Platform",
      Inches(0.5), Inches(2.0), Inches(12.33), Inches(0.88),
      42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(sl, "User & Administrator Manual",
      Inches(0.5), Inches(2.85), Inches(12.33), Inches(0.65),
      26, color=GOLD, align=PP_ALIGN.CENTER)
    R(sl, Inches(3.5), Inches(3.6), Inches(6.33), Inches(0.055), GOLD)

    for i, (label, value) in enumerate([
        ("Version",  "v3.0  |  FastAPI · PostgreSQL · Railway.app"),
        ("Coverage", "All roles: Superadmin · Admin · Cashier"),
        ("Support",  "Contact your Superadmin for account or access issues"),
    ]):
        row_y = Inches(3.8 + i * 0.5)
        T(sl, f"{label}:", Inches(3.5), row_y, Inches(1.8), Inches(0.42),
          13, bold=True, color=GOLD)
        T(sl, value, Inches(5.35), row_y, Inches(7.5), Inches(0.42),
          13, color=RGBColor(0x88, 0xA4, 0xCC))

    T(sl, "እቁብ", Inches(0.5), Inches(5.4), Inches(12.33), Inches(1.6),
      92, bold=True, color=RGBColor(0x0E,0x24,0x42), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════════
s_cover()
s_toc()

divider_slide(1, "Getting Started",     "Login · Dashboard · Navigation",        GREEN)
s_login()
s_dashboard()

divider_slide(2, "Managing Members",    "Add · Edit · Spots · Lifecycle",         BLUE)
s_member_concepts()
s_add_member()
s_member_lifecycle()

divider_slide(3, "Cycles & Draws",      "Create cycle · Run draws · Sell pot",    GOLD)
s_create_cycle()
s_draws()

divider_slide(4, "Payments",            "Record · Batch · Auto-status rules",     RED)
s_payments()
s_payment_rules()

divider_slide(5, "Disbursements",       "Cheques · Guarantors · Vouchers",        PURPLE)
s_disbursements()
s_voucher()

divider_slide(6, "Reports & Analytics", "Financial · Ledger · Export",            CYAN)
s_reports()

divider_slide(7, "SMS Notifications",   "Templates · Auto-send · Log",            GREEN)
s_notifications()
s_sms_templates()

divider_slide(8, "Settings & Admin",    "Roles · Permissions · Security · 2FA",   BLUE)
s_financial_settings()
s_roles()
s_security()

s_glossary()
s_back_cover()

prs.save(OUT)
sys.stdout.buffer.write(b"Done: Equb_Manual.pptx\n")

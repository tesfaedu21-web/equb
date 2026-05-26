"""Generates Equb App Overview PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# Brand colors
DARK_NAVY   = RGBColor(0x0F, 0x24, 0x3E)   # deep navy (brand panel)
GOLD        = RGBColor(0xF5, 0xA6, 0x23)   # Ethiopian gold
GREEN       = RGBColor(0x07, 0x8A, 0x3C)   # Ethiopian green
RED         = RGBColor(0xDC, 0x14, 0x3C)   # Ethiopian red
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF4, 0xF6, 0xF9)
SLATE       = RGBColor(0x3A, 0x4A, 0x5C)
ACCENT_BLUE = RGBColor(0x1A, 0x73, 0xE8)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ── helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_bullet_box(slide, bullets: list[tuple[str, str]], x, y, w, h,
                   title_color=GOLD, body_color=DARK_NAVY, bg=WHITE,
                   icon_color=GREEN):
    """Card with a title row and bullet list."""
    card = add_rect(slide, x, y, w, h, bg)
    # left accent bar
    add_rect(slide, x, y, Inches(0.07), h, icon_color)

    cy = y + Inches(0.18)
    title, *items = bullets
    add_text(slide, title[1], x + Inches(0.18), cy, w - Inches(0.25),
             Inches(0.35), 12, bold=True, color=title_color)
    cy += Inches(0.38)
    for _, line in items:
        add_text(slide, f"• {line}", x + Inches(0.25), cy,
                 w - Inches(0.35), Inches(0.3), 10, color=body_color)
        cy += Inches(0.28)


def add_chip(slide, text, x, y, color: RGBColor):
    chip = add_rect(slide, x, y, Inches(1.9), Inches(0.38), color)
    add_text(slide, text, x, y + Inches(0.05), Inches(1.9), Inches(0.3),
             10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── SLIDE 1 — Title ───────────────────────────────────────────────────────────

def slide_title():
    slide = prs.slides.add_slide(blank_layout)

    # full background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_NAVY)

    # gold ring decorations (thin rectangles as stand-ins)
    for size, x, y in [(Inches(3.5), Inches(9.5), Inches(-1)),
                        (Inches(2.2), Inches(10.8), Inches(4.5))]:
        ring = slide.shapes.add_shape(9, x, y, size, size)  # oval
        ring.fill.background()
        ring.line.color.rgb = GOLD
        ring.line.width = Pt(2)

    # Ethiopian flag stripe accents (left edge)
    add_rect(slide, 0, 0, Inches(0.12), SLIDE_H * 0.33, GREEN)
    add_rect(slide, 0, SLIDE_H * 0.33, Inches(0.12), SLIDE_H * 0.34, GOLD)
    add_rect(slide, 0, SLIDE_H * 0.67, Inches(0.12), SLIDE_H * 0.33, RED)

    # app name
    add_text(slide, "እቁብ", Inches(0.5), Inches(1.2), Inches(6), Inches(1.5),
             80, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

    add_text(slide, "Equb Management Platform",
             Inches(0.5), Inches(2.6), Inches(8), Inches(0.7),
             28, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    # subtitle line
    add_rect(slide, Inches(0.5), Inches(3.25), Inches(3), Inches(0.04), GOLD)

    add_text(slide, "Digital infrastructure for Ethiopia's rotating savings groups",
             Inches(0.5), Inches(3.4), Inches(9), Inches(0.6),
             16, color=RGBColor(0xCC, 0xD6, 0xE8), align=PP_ALIGN.LEFT)

    # stat chips at the bottom
    stats = [("113+ Spots / Cycle", GREEN),
             ("Weekly Draws", ACCENT_BLUE),
             ("SMS Alerts", RED),
             ("Role-Based Access", SLATE)]
    for i, (label, col) in enumerate(stats):
        add_chip(slide, label, Inches(0.5 + i * 2.1), Inches(6.3), col)


# ── SLIDE 2 — What is Equb? ───────────────────────────────────────────────────

def slide_what_is_equb():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)

    # header band
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), DARK_NAVY)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.06), GOLD)
    add_text(slide, "What is Equb?", Inches(0.5), Inches(0.2),
             Inches(10), Inches(0.7), 28, bold=True, color=WHITE)

    # left column — concept
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(5.5), Inches(5.7), WHITE)
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(0.1), Inches(5.7), GREEN)

    add_text(slide, "The Traditional Concept",
             Inches(0.6), Inches(1.55), Inches(5.2), Inches(0.5),
             16, bold=True, color=DARK_NAVY)

    concept_lines = [
        "Equb (እቁብ) is a centuries-old Ethiopian rotating savings and credit group (ROSCA).",
        "",
        "A fixed group of members each contribute an agreed amount every week.",
        "",
        "Each cycle, one member wins the entire pot through a draw.",
        "",
        "Every member eventually receives the full pot — guaranteed.",
        "",
        "It is built on trust, community, and mutual financial support.",
    ]
    y = Inches(2.1)
    for line in concept_lines:
        add_text(slide, line, Inches(0.6), y, Inches(5.1), Inches(0.35),
                 11, color=SLATE)
        y += Inches(0.33)

    # right column — why digitalise
    add_rect(slide, Inches(6.3), Inches(1.35), Inches(6.6), Inches(5.7), WHITE)
    add_rect(slide, Inches(6.3), Inches(1.35), Inches(0.1), Inches(5.7), ACCENT_BLUE)

    add_text(slide, "Why a Digital Platform?",
             Inches(6.5), Inches(1.55), Inches(6.1), Inches(0.5),
             16, bold=True, color=DARK_NAVY)

    why_lines = [
        ("Manual tracking → errors and disputes", RED),
        ("No audit trail for payments and draws", RED),
        ("Hard to scale beyond ~20 members", RED),
        ("SMS reminders impossible at scale", RED),
        ("", WHITE),
        ("Automated payment status tracking", GREEN),
        ("Transparent weekly draw records", GREEN),
        ("113 member spots per cycle", GREEN),
        ("Africa's Talking SMS integration", GREEN),
        ("Role-based admin / cashier access", GREEN),
    ]
    y = Inches(2.1)
    for text, color in why_lines:
        prefix = "✗ " if color == RED else ("✓ " if color == GREEN else "")
        add_text(slide, prefix + text, Inches(6.5), y, Inches(6.1), Inches(0.35),
                 11, color=color if text else WHITE)
        y += Inches(0.33)


# ── SLIDE 3 — Key Features ────────────────────────────────────────────────────

def slide_features():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), DARK_NAVY)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.06), GOLD)
    add_text(slide, "Key Features", Inches(0.5), Inches(0.2),
             Inches(10), Inches(0.7), 28, bold=True, color=WHITE)

    cards = [
        # (title, bullet1, bullet2, bullet3)
        ("Member Management",
         "Track 113+ member spots + 5 association spots",
         "Full & half-share spot splitting",
         "Member status: active / left / stopped paying"),
        ("Cycles & Weekly Draws",
         "Multi-cycle support with per-cycle settings",
         "Automated draw scheduling every week",
         "Status auto-transitions: pending → late → missed"),
        ("Payment Tracking",
         "Weekly payment entry per member",
         "Batch payment recording by cashiers",
         "Payment history with late/missed flags"),
        ("Disbursements",
         "Cheque-based payouts with 3 guarantors required",
         "Voucher deductions tracked per member",
         "End-of-cycle distribution cheques"),
        ("Pot Sales & Transfers",
         "Group week sales, member-to-member transfers",
         "Association spot sales with profit routing",
         "Net pot = gross minus association deduction"),
        ("SMS Notifications",
         "Africa's Talking API integration",
         "6 templates: reminder, winner, missed, on-hold…",
         "Auto-reminders 48 hours before each draw"),
    ]

    cols, rows = 3, 2
    card_w = Inches(4.1)
    card_h = Inches(2.45)
    gap_x  = Inches(0.16)
    gap_y  = Inches(0.18)
    start_x = Inches(0.35)
    start_y = Inches(1.35)

    accent_colors = [GREEN, ACCENT_BLUE, GOLD, RED, GREEN, ACCENT_BLUE]

    for i, (title, b1, b2, b3) in enumerate(cards):
        row, col = divmod(i, cols)
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        add_rect(slide, x, y, card_w, card_h, WHITE)
        add_rect(slide, x, y, Inches(0.08), card_h, accent_colors[i])

        # title
        add_text(slide, title, x + Inches(0.18), y + Inches(0.15),
                 card_w - Inches(0.25), Inches(0.38),
                 13, bold=True, color=DARK_NAVY)
        # divider
        add_rect(slide, x + Inches(0.18), y + Inches(0.58),
                 card_w - Inches(0.35), Inches(0.03), LIGHT_GRAY)
        # bullets
        for j, bullet in enumerate([b1, b2, b3]):
            add_text(slide, f"• {bullet}",
                     x + Inches(0.22), y + Inches(0.72) + j * Inches(0.55),
                     card_w - Inches(0.35), Inches(0.5),
                     10, color=SLATE)


# ── SLIDE 4 — Reports & Analytics ─────────────────────────────────────────────

def slide_reports():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), DARK_NAVY)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.06), GOLD)
    add_text(slide, "Reports & Analytics", Inches(0.5), Inches(0.2),
             Inches(10), Inches(0.7), 28, bold=True, color=WHITE)

    items = [
        ("Financial Summary",       "Total collected vs. expected per cycle; pot size breakdown; association fund balance"),
        ("Payment Status Report",   "Week-by-week member payment grid: paid / late / missed per member"),
        ("Member Status Report",    "Active, exited, and stopped-paying members with join/exit dates"),
        ("Draw History",            "All past draws: winner, week number, pot amount, guarantors"),
        ("Disbursement Ledger",     "Cheque records, voucher deductions, guarantor list per payout"),
        ("Collection Trend Chart",  "Actual ETB collected vs. expected — color-coded weekly bar chart"),
        ("SMS Log",                 "Sent notification history per member: type, timestamp, status"),
        ("Association Expense Log", "Vendor payments and expense categories for the association fund"),
    ]

    y = Inches(1.35)
    for title, desc in items:
        add_rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.58), WHITE)
        add_rect(slide, Inches(0.4), y, Inches(0.08), Inches(0.58), ACCENT_BLUE)
        add_text(slide, title, Inches(0.6), y + Inches(0.08),
                 Inches(3.2), Inches(0.42), 11, bold=True, color=DARK_NAVY)
        add_text(slide, desc, Inches(3.9), y + Inches(0.08),
                 Inches(9.0), Inches(0.42), 10, color=SLATE)
        y += Inches(0.65)


# ── SLIDE 5 — Tech Stack ──────────────────────────────────────────────────────

def slide_tech():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_NAVY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), RGBColor(0x06, 0x14, 0x28))
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.06), GOLD)
    add_text(slide, "Technology Stack", Inches(0.5), Inches(0.2),
             Inches(10), Inches(0.7), 28, bold=True, color=WHITE)

    layers = [
        ("Frontend",  ACCENT_BLUE,
         ["Jinja2 HTML templates", "Tailwind CSS", "Chart.js visualisations", "Vanilla JavaScript"]),
        ("Backend",   GREEN,
         ["FastAPI (Python)", "Uvicorn ASGI server", "APScheduler (cron jobs)", "SQLAlchemy ORM"]),
        ("Database",  GOLD,
         ["PostgreSQL (production)", "SQLite (local dev)", "Railway-hosted Postgres", "Alembic migrations"]),
        ("Services",  RED,
         ["Africa's Talking SMS API", "Railway.app hosting", "PBKDF2 password hashing", "CSP / HTTPS / CSRF"]),
    ]

    col_w = Inches(3.1)
    for i, (label, color, items) in enumerate(layers):
        x = Inches(0.35) + i * (col_w + Inches(0.15))
        y = Inches(1.35)

        # header chip
        add_rect(slide, x, y, col_w, Inches(0.5), color)
        add_text(slide, label, x, y + Inches(0.07), col_w, Inches(0.36),
                 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # items
        card_h = Inches(0.72)
        for j, item in enumerate(items):
            iy = y + Inches(0.6) + j * (card_h + Inches(0.08))
            add_rect(slide, x, iy, col_w, card_h, RGBColor(0x1A, 0x2E, 0x4A))
            add_rect(slide, x, iy, Inches(0.07), card_h, color)
            add_text(slide, item, x + Inches(0.18), iy + Inches(0.18),
                     col_w - Inches(0.25), Inches(0.36), 12, color=WHITE)


# ── SLIDE 6 — Security & Roles ────────────────────────────────────────────────

def slide_security():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), DARK_NAVY)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.06), GOLD)
    add_text(slide, "Security & Access Control", Inches(0.5), Inches(0.2),
             Inches(10), Inches(0.7), 28, bold=True, color=WHITE)

    # left: roles
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(5.8), Inches(5.7), WHITE)
    add_rect(slide, Inches(0.4), Inches(1.35), Inches(0.1), Inches(5.7), GOLD)
    add_text(slide, "Role-Based Access", Inches(0.6), Inches(1.55),
             Inches(5.4), Inches(0.5), 16, bold=True, color=DARK_NAVY)

    roles = [
        ("Superadmin", "Full access: settings, users, all data, reports"),
        ("Admin",      "Cycle management, draws, disbursements, reports"),
        ("Cashier",    "Payment entry and member viewing only"),
    ]
    y = Inches(2.2)
    for role, desc in roles:
        add_rect(slide, Inches(0.6), y, Inches(5.4), Inches(1.1), LIGHT_GRAY)
        add_rect(slide, Inches(0.6), y, Inches(0.08), Inches(1.1), ACCENT_BLUE)
        add_text(slide, role, Inches(0.8), y + Inches(0.1),
                 Inches(4.9), Inches(0.4), 13, bold=True, color=DARK_NAVY)
        add_text(slide, desc, Inches(0.8), y + Inches(0.52),
                 Inches(4.9), Inches(0.45), 10, color=SLATE)
        y += Inches(1.25)

    # right: security measures
    add_rect(slide, Inches(6.6), Inches(1.35), Inches(6.35), Inches(5.7), WHITE)
    add_rect(slide, Inches(6.6), Inches(1.35), Inches(0.1), Inches(5.7), RED)
    add_text(slide, "Security Measures", Inches(6.8), Inches(1.55),
             Inches(6.0), Inches(0.5), 16, bold=True, color=DARK_NAVY)

    measures = [
        "PBKDF2 password hashing (no plain-text storage)",
        "Rate limiting: 5 failed logins → 5-minute lockout",
        "CSRF protection on all state-changing requests",
        "Content Security Policy (CSP) headers enforced",
        "HTTPS enforcement in production",
        "Session middleware with secure cookie flags",
        "Granular per-role permission matrix",
        "Audit trail for payments, draws, and disbursements",
    ]
    y = Inches(2.15)
    for m in measures:
        add_text(slide, f"✓  {m}", Inches(6.8), y, Inches(5.9), Inches(0.45),
                 11, color=SLATE)
        y += Inches(0.58)


# ── SLIDE 7 — Architecture Diagram ───────────────────────────────────────────

def slide_architecture():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF0, 0xF2, 0xF5))
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), DARK_NAVY)
    add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.06), GOLD)
    add_text(slide, "System Architecture", Inches(0.5), Inches(0.2),
             Inches(10), Inches(0.7), 28, bold=True, color=WHITE)

    # Browser
    add_rect(slide, Inches(0.4), Inches(1.8), Inches(2.2), Inches(0.8), ACCENT_BLUE)
    add_text(slide, "Browser\n(Admin / Cashier)", Inches(0.4), Inches(1.88),
             Inches(2.2), Inches(0.65), 10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # arrow → FastAPI
    add_rect(slide, Inches(2.62), Inches(2.12), Inches(1.1), Inches(0.06), SLATE)

    # FastAPI box
    add_rect(slide, Inches(3.75), Inches(1.6), Inches(3.0), Inches(1.2), DARK_NAVY)
    add_text(slide, "FastAPI Server\n(Python / Uvicorn)",
             Inches(3.75), Inches(1.75), Inches(3.0), Inches(0.9),
             12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Routers below FastAPI
    routers = ["auth", "members", "draws", "payments", "reports", "disbursements", "notifications", "settings"]
    rx = Inches(3.2)
    for i, r in enumerate(routers):
        col = i % 4
        row = i // 4
        bx = Inches(3.2) + col * Inches(1.72)
        by = Inches(3.1) + row * Inches(0.62)
        add_rect(slide, bx, by, Inches(1.6), Inches(0.48), RGBColor(0x1A, 0x2E, 0x4A))
        add_text(slide, r, bx, by + Inches(0.08), Inches(1.6), Inches(0.32),
                 9, color=WHITE, align=PP_ALIGN.CENTER)

    # arrow to DB
    add_rect(slide, Inches(7.3), Inches(2.12), Inches(0.8), Inches(0.06), SLATE)

    # Postgres
    add_rect(slide, Inches(8.15), Inches(1.6), Inches(2.3), Inches(1.2), GREEN)
    add_text(slide, "PostgreSQL\n(Railway)", Inches(8.15), Inches(1.78),
             Inches(2.3), Inches(0.85), 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # SMS
    add_rect(slide, Inches(8.15), Inches(3.2), Inches(2.3), Inches(1.0), GOLD)
    add_text(slide, "Africa's Talking\nSMS API", Inches(8.15), Inches(3.35),
             Inches(2.3), Inches(0.7), 11, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)

    # APScheduler
    add_rect(slide, Inches(3.75), Inches(5.0), Inches(3.0), Inches(0.9), RED)
    add_text(slide, "APScheduler\n(Nightly jobs + pre-draw reminders)",
             Inches(3.75), Inches(5.1), Inches(3.0), Inches(0.7),
             10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Railway hosting banner
    add_rect(slide, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.7), DARK_NAVY)
    add_text(slide, "Hosted on Railway.app  |  Auto-deploy from Git  |  Managed Postgres add-on",
             Inches(0.4), Inches(6.38), Inches(12.5), Inches(0.55),
             11, color=GOLD, align=PP_ALIGN.CENTER)


# ── SLIDE 8 — Roadmap / Closing ───────────────────────────────────────────────

def slide_closing():
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_NAVY)

    # decorative rings
    for size, x, y in [(Inches(4), Inches(-1.2), Inches(-1)),
                        (Inches(2.5), Inches(10.5), Inches(5))]:
        ring = slide.shapes.add_shape(9, x, y, size, size)
        ring.fill.background()
        ring.line.color.rgb = RGBColor(0x2A, 0x4A, 0x6A)
        ring.line.width = Pt(3)

    # Ethiopian flag stripes (right edge)
    add_rect(slide, SLIDE_W - Inches(0.12), 0, Inches(0.12), SLIDE_H * 0.33, GREEN)
    add_rect(slide, SLIDE_W - Inches(0.12), SLIDE_H * 0.33, Inches(0.12), SLIDE_H * 0.34, GOLD)
    add_rect(slide, SLIDE_W - Inches(0.12), SLIDE_H * 0.67, Inches(0.12), SLIDE_H * 0.33, RED)

    add_text(slide, "Built for Ethiopia's",
             Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             36, color=RGBColor(0xCC, 0xD6, 0xE8), align=PP_ALIGN.CENTER)
    add_text(slide, "Financial Communities",
             Inches(1), Inches(1.9), Inches(11), Inches(1.0),
             48, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(4.5), Inches(3.0), Inches(4.3), Inches(0.04), GOLD)

    pillars = ["Transparent", "Automated", "Scalable", "Secure"]
    for i, p in enumerate(pillars):
        bx = Inches(0.8) + i * Inches(3.0)
        add_rect(slide, bx, Inches(3.3), Inches(2.6), Inches(0.6),
                 [GREEN, ACCENT_BLUE, GOLD, RED][i])
        add_text(slide, p, bx, Inches(3.38), Inches(2.6), Inches(0.44),
                 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Thank you",
             Inches(1), Inches(4.5), Inches(11), Inches(0.7),
             32, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "እቁብ — Equb Management Platform",
             Inches(1), Inches(5.2), Inches(11), Inches(0.5),
             14, color=RGBColor(0x88, 0xA0, 0xBB), align=PP_ALIGN.CENTER)


# ── Build & save ──────────────────────────────────────────────────────────────

slide_title()
slide_what_is_equb()
slide_features()
slide_reports()
slide_tech()
slide_security()
slide_architecture()
slide_closing()

out_path = os.path.join(os.path.dirname(__file__), "Equb_Presentation.pptx")
prs.save(out_path)
import sys
sys.stdout.buffer.write(("Saved: " + out_path + "\n").encode("utf-8"))
